from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Slide 1 – Title
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Internship on Django Framework"
subtitle.text = "Using Python for Backend and Frontend Development"

# Slide 2 – Introduction
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Introduction"
content.text = (
    "This internship provided hands-on experience in Django framework using Python.\n"
    "I learned how to design backend systems, integrate frontend templates, "
    "and connect applications with databases."
)

# Slide 3 – Objective
slide = prs.slides.add_slide(slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Objective of Internship Training"
content.text = (
    "• To gain practical knowledge of Django framework.\n"
    "• To understand backend development with Python.\n"
    "• To explore frontend integration using HTML, CSS, and JavaScript.\n"
    "• To learn how to connect web applications with databases (PostgreSQL/MySQL).\n"
    "• To enhance problem-solving and real-time project development skills."
)

# Slide 4 – Tools and Technology
slide = prs.slides.add_slide(slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Tools and Technology Used"
content.text = (
    "• Python\n"
    "• Django Framework\n"
    "• HTML, CSS, JavaScript\n"
    "• PostgreSQL / MySQL\n"
    "• Visual Studio Code / PyCharm\n"
    "• Git & GitHub for Version Control"
)

# Slide 5 – Overview of Training
slide = prs.slides.add_slide(slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Overview of Training Work Undertaken"
content.text = (
    "• Developed login and signup system.\n"
    "• Implemented database connectivity.\n"
    "• Created CRUD operations using Django models.\n"
    "• Designed templates for frontend integration.\n"
    "• Built modules for food ordering system.\n"
    "• Learned deployment basics."
)

# Slide 6 – Description & Screenshots
slide = prs.slides.add_slide(slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Description & Screenshots"
content.text = (
    "Screenshots of my work:\n"
    "• Login Page\n"
    "• Signup Page\n"
    "• Food Ordering Module\n"
    "• Cart & Order Summary"
)

# Slide 7 – Conclusion
slide = prs.slides.add_slide(slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Conclusion"
content.text = (
    "• The internship helped me strengthen my skills in Python and Django.\n"
    "• I gained practical exposure to backend and frontend integration.\n"
    "• Learned how to manage databases with PostgreSQL.\n"
    "• Improved my understanding of web application development lifecycle."
)

# Slide 8 – References
slide = prs.slides.add_slide(slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "References"
content.text = (
    "• Django Official Documentation (https://www.djangoproject.com/)\n"
    "• Python Documentation (https://docs.python.org/)\n"
    "• W3Schools (HTML, CSS, JavaScript)\n"
    "• Stack Overflow & GitHub Community"
)

# Save the presentation
prs.save("Django_Internship_Presentation.pptx")
print("Presentation created successfully: Django_Internship_Presentation.pptx")
